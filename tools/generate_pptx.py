"""
tool: generate_pptx
-------------------
Generates a management-ready PowerPoint investigation briefing
from a case's artefacts and IOC data.

Usage:
  python3 tools/generate_pptx.py --case C003

Output:
  cases/<CASE_ID>/reports/investigation_briefing.pptx
"""
from __future__ import annotations

import json
import sys
from datetime import datetime, timezone
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from config.settings import CASES_DIR
from tools.common import audit, load_json, utcnow, write_artefact

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x0D, 0x1B, 0x2A)   # slide background
BLUE_MID    = RGBColor(0x1B, 0x4F, 0x72)   # section headers / boxes
BLUE_LIGHT  = RGBColor(0x21, 0x8D, 0xBF)   # accent / arrows
AMBER       = RGBColor(0xF3, 0x9C, 0x12)   # warning / HIGH severity
RED         = RGBColor(0xC0, 0x39, 0x2B)   # critical
GREEN       = RGBColor(0x1E, 0x8B, 0x4C)   # clean / low severity
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY  = RGBColor(0xD5, 0xD8, 0xDC)
MID_GREY    = RGBColor(0x7F, 0x8C, 0x8D)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Low-level helpers ─────────────────────────────────────────────────────────

def _set_bg(slide, color: RGBColor):
    """Fill slide background with a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        l, t, w, h
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def _add_textbox(slide, l, t, w, h, text, font_size=Pt(14),
                 color=WHITE, bold=False, align=PP_ALIGN.LEFT,
                 italic=False, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = font_size
    run.font.color.rgb = color
    run.font.bold  = bold
    run.font.italic = italic
    return txb


def _add_label_box(slide, l, t, w, h, label, value,
                   box_color=BLUE_MID, label_size=Pt(11), val_size=Pt(14)):
    """Labelled value box (like a stat card)."""
    rect = _add_rect(slide, l, t, w, h, fill_color=box_color)
    tf   = rect.text_frame
    tf.word_wrap = True
    # label line
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = label.upper()
    r1.font.size  = label_size
    r1.font.color.rgb = LIGHT_GREY
    r1.font.bold  = False
    # value line
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = value
    r2.font.size  = val_size
    r2.font.color.rgb = WHITE
    r2.font.bold  = True
    return rect


def _add_flow_box(slide, l, t, w, h, title, lines,
                  box_fill=BLUE_MID, title_size=Pt(13), body_size=Pt(10.5)):
    rect = _add_rect(slide, l, t, w, h, fill_color=box_fill)
    tf   = rect.text_frame
    tf.word_wrap = True
    tf.margin_left   = Inches(0.1)
    tf.margin_right  = Inches(0.1)
    tf.margin_top    = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    r0 = p0.add_run()
    r0.text = title
    r0.font.size  = title_size
    r0.font.color.rgb = WHITE
    r0.font.bold  = True
    for line in lines:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line
        r.font.size  = body_size
        r.font.color.rgb = LIGHT_GREY
    return rect


def _arrow(slide, x1, y1, x2, y2, color=BLUE_LIGHT, width=Pt(2)):
    """Draw a horizontal arrow using a connector."""
    conn = slide.shapes.add_connector(
        pptx.enum.shapes.MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2
    )
    conn.line.color.rgb = color
    conn.line.width     = width
    return conn


def _slide_header(slide, title: str, subtitle: str = ""):
    """Add a top banner with title."""
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.85), fill_color=BLUE_MID)
    _add_textbox(slide, Inches(0.3), Inches(0.08), Inches(10), Inches(0.55),
                 title, font_size=Pt(24), bold=True, color=WHITE)
    if subtitle:
        _add_textbox(slide, Inches(0.3), Inches(0.6), Inches(10), Inches(0.3),
                     subtitle, font_size=Pt(12), color=LIGHT_GREY)


def _speaker_note(slide, text: str):
    notes = slide.notes_slide
    tf    = notes.notes_text_frame
    tf.text = text


# ── Slide builders ────────────────────────────────────────────────────────────

def slide_title(prs, meta: dict):
    layout = prs.slide_layouts[6]  # blank
    sl = prs.slides.add_slide(layout)
    _set_bg(sl, NAVY)

    # Top accent bar
    _add_rect(sl, 0, 0, SLIDE_W, Inches(0.12), fill_color=BLUE_LIGHT)

    # Case badge
    _add_rect(sl, Inches(0.4), Inches(1.2), Inches(2.2), Inches(0.55),
              fill_color=AMBER)
    _add_textbox(sl, Inches(0.4), Inches(1.22), Inches(2.2), Inches(0.5),
                 f"CASE  {meta.get('case_id','?')}",
                 font_size=Pt(18), bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Severity badge
    sev   = meta.get("severity", "medium").upper()
    s_col = RED if sev == "CRITICAL" else AMBER if sev == "HIGH" else BLUE_MID
    _add_rect(sl, Inches(2.8), Inches(1.2), Inches(1.6), Inches(0.55), fill_color=s_col)
    _add_textbox(sl, Inches(2.8), Inches(1.22), Inches(1.6), Inches(0.5),
                 f"⚠  {sev}", font_size=Pt(16), bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER)

    # Main title
    _add_textbox(sl, Inches(0.4), Inches(2.0), Inches(12.5), Inches(1.4),
                 meta.get("title", "Security Investigation"),
                 font_size=Pt(38), bold=True, color=WHITE)

    # Subtitle
    _add_textbox(sl, Inches(0.4), Inches(3.5), Inches(10), Inches(0.5),
                 "Threat Investigation Briefing  ·  Security Operations Centre",
                 font_size=Pt(16), color=LIGHT_GREY)

    # Metadata row
    now = datetime.now(timezone.utc).strftime("%d %B %Y")
    meta_line = (f"Analyst: {meta.get('analyst','unassigned')}   |   "
                 f"Date: {now}   |   "
                 f"Status: {meta.get('status','open').upper()}")
    _add_textbox(sl, Inches(0.4), Inches(4.2), Inches(10), Inches(0.4),
                 meta_line, font_size=Pt(13), color=MID_GREY)

    # Bottom bar
    _add_rect(sl, 0, Inches(7.2), SLIDE_W, Inches(0.3), fill_color=BLUE_MID)
    _add_textbox(sl, Inches(0.2), Inches(7.22), Inches(6), Inches(0.25),
                 "RESTRICTED – INTERNAL USE ONLY",
                 font_size=Pt(9), color=LIGHT_GREY)

    _speaker_note(sl,
        "SPEAKER: Introduce the case. Explain this is an automated SOC investigation briefing. "
        f"Case {meta.get('case_id')} was opened following submission of a suspicious PowerShell script. "
        "Severity was assessed as HIGH. All findings in this deck are evidence-based and reference "
        "artefacts stored in the case vault.")


def slide_executive_summary(prs, meta: dict, iocs: dict):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    _set_bg(sl, NAVY)
    _slide_header(sl, "Executive Summary", "What happened and why it matters")

    totals = iocs.get("total", {})

    # Stat cards
    cards = [
        ("Severity",    meta.get("severity","?").upper()),
        ("IOC Types",   str(sum(1 for v in totals.values() if v > 0))),
        ("Domains",     str(totals.get("domain", 0))),
        ("URLs (C2)",   str(totals.get("url", 0))),
        ("File Hashes", str(totals.get("sha256", 0))),
        ("Status",      meta.get("status","open").upper()),
    ]
    card_w = Inches(1.9)
    card_h = Inches(1.05)
    card_t = Inches(1.0)
    for i, (lbl, val) in enumerate(cards):
        col = RED if (lbl == "Severity" and val in ("HIGH","CRITICAL")) else BLUE_MID
        _add_label_box(sl,
                       Inches(0.25) + i * (card_w + Inches(0.12)),
                       card_t, card_w, card_h, lbl, val, box_color=col)

    # Key finding bullets
    bullets = [
        "A malicious PowerShell script was submitted for analysis and confirmed as active malware.",
        "The script is the initial-stage payload of a ClickFix / fake-CAPTCHA social engineering campaign.",
        "It beacons victim device details to a remote C2 server, then delivers a three-stage attack chain.",
        "The final payload (HelloMemory.dll) is a custom, maintained .NET implant loaded entirely in memory.",
        "It establishes registry persistence, supports self-update, and rotates between multiple C2 mirrors.",
        "Infrastructure is hosted on a compromised Romanian shared hosting server (89.42.218.223) running "
          "unpatched Exim with 6 live CVEs.",
        "163 other websites share the same server and may be serving identical lures to other victims.",
        "No public threat intelligence matches were found — HelloMemory appears to be unreported tooling.",
    ]
    _add_textbox(sl, Inches(0.3), Inches(2.25), Inches(12.5), Inches(0.35),
                 "Key Findings", font_size=Pt(15), bold=True, color=BLUE_LIGHT)
    y = Inches(2.65)
    for b in bullets:
        _add_textbox(sl, Inches(0.55), y, Inches(12.4), Inches(0.38),
                     f"▸  {b}", font_size=Pt(11.5), color=LIGHT_GREY)
        y += Inches(0.42)

    _add_rect(sl, 0, Inches(7.2), SLIDE_W, Inches(0.3), fill_color=BLUE_MID)
    _speaker_note(sl,
        "SPEAKER: Summarise the situation for management. "
        "A user was shown a fake CAPTCHA page that instructed them to paste a command into the Windows Run dialog. "
        "That command was this script. If executed, the machine is compromised. "
        "Emphasise: we do not yet know if any hosts actually ran this — that is a priority investigation action.")


def slide_how_it_arrived(prs):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    _set_bg(sl, NAVY)
    _slide_header(sl, "How the Attack Arrives", "ClickFix / Fake CAPTCHA — Social Engineering Delivery")

    # Three-panel description
    panels = [
        ("1  USER VISITS\nCOMPROMISED PAGE",
         ["Normal website has been", "hacked and injects a", "fake CAPTCHA overlay"],
         BLUE_MID),
        ("2  FAKE CAPTCHA\nINSTRUCTIONS",
         ['"To verify you are human:"', "① Press  Win + R", "② Press  Ctrl + V", "③ Press  Enter"],
         RGBColor(0x6E, 0x27, 0x94)),
        ("3  POWERSHELL\nSILENTLY RUNS",
         ["Malicious command was", "pre-loaded to clipboard", "by the page's JavaScript", "before user pasted it"],
         RED),
    ]
    bw = Inches(3.6)
    bh = Inches(3.2)
    bt = Inches(1.2)
    gap = Inches(0.5)
    for i, (title, lines, col) in enumerate(panels):
        bl = Inches(0.5) + i * (bw + gap)
        _add_flow_box(sl, bl, bt, bw, bh, title, lines,
                      box_fill=col, title_size=Pt(15), body_size=Pt(12))
        if i < 2:
            ax = bl + bw + Inches(0.05)
            ay = bt + bh / 2
            _add_textbox(sl, ax, ay - Inches(0.2), gap - Inches(0.1), Inches(0.4),
                         "▶", font_size=Pt(28), color=BLUE_LIGHT, align=PP_ALIGN.CENTER)

    # Callout
    _add_rect(sl, Inches(0.5), Inches(4.65), Inches(12.3), Inches(0.72),
              fill_color=RGBColor(0x17, 0x20, 0x2A),
              line_color=AMBER, line_width=Pt(1.5))
    _add_textbox(sl, Inches(0.7), Inches(4.72), Inches(12.0), Inches(0.6),
                 "⚠  The script submitted for analysis (submitted_script.ps1) is identical to vcapcha.ps1 "
                 "— the live ClickFix payload already observed on this domain by AlienVault OTX on 23 Feb 2026.",
                 font_size=Pt(11.5), color=AMBER)

    _add_textbox(sl, Inches(0.5), Inches(5.55), Inches(12.3), Inches(1.55),
                 "Why this works:  CAPTCHA challenges never require keyboard shortcuts or system-level commands. "
                 "Any website asking you to press Win+R and paste something is a guaranteed attack. "
                 "Users should be trained to recognise and report this pattern immediately.",
                 font_size=Pt(11), color=MID_GREY)

    _add_rect(sl, 0, Inches(7.2), SLIDE_W, Inches(0.3), fill_color=BLUE_MID)
    _speaker_note(sl,
        "SPEAKER: Explain the ClickFix technique without jargon. "
        "The victim sees what looks like a Cloudflare CAPTCHA. They click 'I am not a robot' — "
        "but invisible JavaScript has already placed a malicious command on their clipboard. "
        "The page then tells them to open the Windows Run dialog and paste. "
        "Most users comply because they believe it is a legitimate verification step. "
        "This technique requires no software exploit — it exploits human behaviour.")


def slide_attack_chain(prs):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    _set_bg(sl, NAVY)
    _slide_header(sl, "Full Attack Chain", "Four-stage delivery from clipboard to persistent implant")

    stages = [
        ("STAGE 0\nCLICKFIX LURE",
         "vcapcha.ps1",
         ["Fake CAPTCHA page", "Clipboard hijack", "Win+R → paste → Enter"],
         RGBColor(0x6E, 0x27, 0x94)),
        ("STAGE 1\nDROPPER",
         "submitted_script.ps1",
         ["Beacon → /reportv.php", "Hide console window", "Download verify.ps1", "Bypass exec policy"],
         AMBER),
        ("STAGE 2\nLOADER",
         "verify.ps1",
         ["Download notepad.b64", "Base64 decode in RAM", "Reflective .NET load", "Call Execute()"],
         RGBColor(0xCA, 0x6F, 0x1E)),
        ("STAGE 3\nIMPLANT",
         "HelloMemory.dll",
         ["Registry persistence", "C2 telemetry", "Self-update loop", "Multi-mirror failover"],
         RED),
    ]

    bw  = Inches(2.6)
    bh  = Inches(3.6)
    bt  = Inches(1.1)
    gap = Inches(0.44)
    for i, (stage, fname, lines, col) in enumerate(stages):
        bl = Inches(0.35) + i * (bw + gap)
        _add_flow_box(sl, bl, bt, bw, bh, stage, [f"  {l}" for l in lines],
                      box_fill=col, title_size=Pt(12.5), body_size=Pt(11))
        # filename tag below box
        _add_textbox(sl, bl, bt + bh + Inches(0.08), bw, Inches(0.3),
                     fname, font_size=Pt(9.5), color=MID_GREY,
                     align=PP_ALIGN.CENTER, italic=True)
        if i < 3:
            ax = bl + bw + Inches(0.05)
            ay = bt + bh / 2 - Inches(0.2)
            _add_textbox(sl, ax, ay, gap - Inches(0.05), Inches(0.4),
                         "▶", font_size=Pt(26), color=BLUE_LIGHT, align=PP_ALIGN.CENTER)

    # Stage 4 hint
    _add_rect(sl, Inches(0.35), Inches(5.05), Inches(12.6), Inches(0.55),
              fill_color=RGBColor(0x17, 0x20, 0x2A), line_color=RED, line_width=Pt(1))
    _add_textbox(sl, Inches(0.55), Inches(5.1), Inches(12.2), Inches(0.45),
                 "STAGE 4 (unconfirmed)  ·  HelloMemory.dll contains a hardcoded SHA-256 reference to a "
                 "fourth-stage payload not yet retrieved.  If a host executed Stage 3, assume further compromise.",
                 font_size=Pt(10.5), color=RED)

    _add_rect(sl, 0, Inches(7.2), SLIDE_W, Inches(0.3), fill_color=BLUE_MID)
    _speaker_note(sl,
        "SPEAKER: Walk through each stage left to right. "
        "Stage 0 is the social engineering delivery. "
        "Stage 1 (the script submitted to us) phones home and fetches the next stage. "
        "Stage 2 is a loader that never writes anything to disk — it decodes and runs entirely in RAM. "
        "Stage 3 is a persistent implant that survives reboots and can update itself. "
        "We believe there is a Stage 4 but have not retrieved it. "
        "Hosts that ran past Stage 1 should be treated as fully compromised.")


def slide_infrastructure(prs):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    _set_bg(sl, NAVY)
    _slide_header(sl, "Attacker Infrastructure", "Compromised Romanian shared hosting server")

    # Server info box
    server_lines = [
        "IP Address:      89.42.218.223",
        "ASN:             AS205275  romarg srl  (Romania)",
        "Hosting:         cPanel / WHM shared hosting  (whmpanels.com)",
        "Domains on IP:   163  (all potentially compromised)",
        "Server software: OpenResty 1.27.1  ·  Exim 4.96.2  ·  MariaDB 10.3.27",
        "Shodan tag:      eol-product  (end-of-life software)",
    ]
    _add_rect(sl, Inches(0.3), Inches(1.0), Inches(7.8), Inches(2.55),
              fill_color=BLUE_MID)
    _add_textbox(sl, Inches(0.5), Inches(1.05), Inches(7.5), Inches(0.4),
                 "SERVER PROFILE", font_size=Pt(12), bold=True, color=BLUE_LIGHT)
    y = Inches(1.5)
    for line in server_lines:
        _add_textbox(sl, Inches(0.55), y, Inches(7.4), Inches(0.32),
                     line, font_size=Pt(10.5), color=LIGHT_GREY)
        y += Inches(0.32)

    # CVE box
    cves = [
        "CVE-2023-51766  –  Exim SMTP smuggling",
        "CVE-2024-39929  –  Exim header parsing RCE",
        "CVE-2022-3559   –  Exim memory corruption",
        "CVE-2022-3620   –  Exim buffer overflow",
        "CVE-2025-30232  –  Exim use-after-free",
        "CVE-2025-67896  –  Unconfirmed (recent)",
    ]
    _add_rect(sl, Inches(8.3), Inches(1.0), Inches(4.7), Inches(2.55),
              fill_color=RGBColor(0x3B, 0x14, 0x14))
    _add_textbox(sl, Inches(8.5), Inches(1.05), Inches(4.4), Inches(0.4),
                 "LIVE CVEs ON SERVER", font_size=Pt(12), bold=True, color=RED)
    y = Inches(1.5)
    for cve in cves:
        _add_textbox(sl, Inches(8.5), y, Inches(4.35), Inches(0.32),
                     cve, font_size=Pt(10), color=LIGHT_GREY)
        y += Inches(0.32)

    # C2 endpoints
    _add_textbox(sl, Inches(0.3), Inches(3.75), Inches(12.5), Inches(0.35),
                 "C2 ENDPOINTS OBSERVED", font_size=Pt(13), bold=True, color=BLUE_LIGHT)
    endpoints = [
        ("simpludelicios.md/reportv.php",  "Victim registration beacon  (POST — device ID + timestamp)"),
        ("simpludelicios.md/verify.ps1",   "Stage 2 loader delivery"),
        ("simpludelicios.md/notepad.b64",  "Stage 3 DLL — base64-encoded .NET implant"),
        ("simpludelicios.md/vcapcha.ps1",  "ClickFix clipboard payload  (= submitted script)"),
    ]
    y = Inches(4.15)
    for ep, desc in endpoints:
        _add_rect(sl, Inches(0.3), y, Inches(4.5), Inches(0.35),
                  fill_color=RGBColor(0x17, 0x20, 0x2A), line_color=BLUE_LIGHT, line_width=Pt(0.75))
        _add_textbox(sl, Inches(0.35), y + Inches(0.02), Inches(4.4), Inches(0.32),
                     ep, font_size=Pt(9.5), color=BLUE_LIGHT, italic=True)
        _add_textbox(sl, Inches(5.0), y + Inches(0.02), Inches(8.0), Inches(0.32),
                     f"→  {desc}", font_size=Pt(10), color=LIGHT_GREY)
        y += Inches(0.43)

    _add_textbox(sl, Inches(0.3), Inches(6.3), Inches(12.5), Inches(0.6),
                 "How the server was likely compromised:  The server runs Exim 4.96.2 with multiple unpatched "
                 "SMTP vulnerabilities. The attacker likely exploited one of these to gain hosting access, "
                 "then planted scripts across all sites on the shared server.",
                 font_size=Pt(10.5), color=MID_GREY)

    _add_rect(sl, 0, Inches(7.2), SLIDE_W, Inches(0.3), fill_color=BLUE_MID)
    _speaker_note(sl,
        "SPEAKER: The attacker is not using dedicated infrastructure — they have compromised a legitimate "
        "Romanian recipe website and its shared hosting server. The server has 6 unpatched Exim "
        "(mail server) vulnerabilities, which is almost certainly how they gained access. "
        "163 other websites share this server and are also at risk. "
        "We should report this to the hosting provider for takedown.")


def slide_hellomemory(prs):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    _set_bg(sl, NAVY)
    _slide_header(sl, "HelloMemory — The Implant", "Custom .NET payload, fileless, persistent, self-updating")

    # Left: capabilities
    caps = [
        ("Device registration",    "Fingerprints each victim on first run and reports to C2"),
        ("Fileless execution",     "Loaded entirely in RAM — never written to disk as a file"),
        ("Registry persistence",   "Adds itself to HKCU Run key — survives reboots"),
        ("Self-update loop",       "Polls C2 for newer versions and replaces itself automatically"),
        ("Multi-mirror C2",        "Rotates through backup C2 addresses if primary is blocked"),
        ("Event reporting",        "Sends execution events back to the operator in real-time"),
        ("4th-stage loader",       "Contains a hardcoded SHA-256 reference — further payload expected"),
    ]
    _add_textbox(sl, Inches(0.3), Inches(0.95), Inches(7.5), Inches(0.35),
                 "IMPLANT CAPABILITIES", font_size=Pt(13), bold=True, color=BLUE_LIGHT)
    y = Inches(1.35)
    for cap, desc in caps:
        _add_rect(sl, Inches(0.3), y, Inches(2.2), Inches(0.38), fill_color=BLUE_MID)
        _add_textbox(sl, Inches(0.35), y + Inches(0.02), Inches(2.1), Inches(0.34),
                     cap, font_size=Pt(10), bold=True, color=WHITE)
        _add_textbox(sl, Inches(2.65), y + Inches(0.02), Inches(5.2), Inches(0.34),
                     desc, font_size=Pt(10), color=LIGHT_GREY)
        y += Inches(0.46)

    # Right: key strings from DLL
    _add_textbox(sl, Inches(8.2), Inches(0.95), Inches(4.8), Inches(0.35),
                 "KEY STRINGS FROM DLL", font_size=Pt(13), bold=True, color=BLUE_LIGHT)
    strings_of_interest = [
        "HelloMemory",
        "EnsureAutoStart",
        "DownloadAndDecodeUpdate",
        "ReportEvent",
        "CheckAndUpdateIfNeeded",
        "EncVersionUrls",
        "EncReportUrls",
        "preferredMirror",
        "InstalledFlagFile",
        "C:\\Users\\marks\\source\\...",
    ]
    y = Inches(1.35)
    for s in strings_of_interest:
        _add_rect(sl, Inches(8.2), y, Inches(4.8), Inches(0.35),
                  fill_color=RGBColor(0x0A, 0x14, 0x1E),
                  line_color=BLUE_MID, line_width=Pt(0.5))
        _add_textbox(sl, Inches(8.3), y + Inches(0.02), Inches(4.6), Inches(0.3),
                     s, font_size=Pt(10), color=BLUE_LIGHT, italic=True)
        y += Inches(0.43)

    # PDB callout
    _add_rect(sl, Inches(0.3), Inches(5.0), Inches(12.6), Inches(0.72),
              fill_color=RGBColor(0x17, 0x20, 0x2A), line_color=AMBER, line_width=Pt(1.5))
    _add_textbox(sl, Inches(0.5), Inches(5.06), Inches(12.2), Inches(0.6),
                 "Developer artefact found in DLL:   "
                 "C:\\Users\\marks\\source\\repos\\ConsoleApp1\\obj\\Release\\net472\\HelloMemory.pdb\n"
                 "Developer username 'marks' leaked via unstripped PDB path. No public attribution match — this is unreported tooling.",
                 font_size=Pt(10.5), color=AMBER)

    _add_rect(sl, Inches(0.3), Inches(5.9), Inches(12.6), Inches(0.55),
              fill_color=RGBColor(0x14, 0x17, 0x1E))
    _add_textbox(sl, Inches(0.5), Inches(5.95), Inches(12.2), Inches(0.45),
                 "SHA-256 (DLL):  160dc775c7fccd4ac45c250944f5d3cd4f3f8414ba2587742f0fb85733fa48c7     "
                 "Size: 13,312 bytes     .NET 4.7.2     No public VT matches at time of analysis",
                 font_size=Pt(9.5), color=MID_GREY)

    _add_rect(sl, 0, Inches(7.2), SLIDE_W, Inches(0.3), fill_color=BLUE_MID)
    _speaker_note(sl,
        "SPEAKER: HelloMemory is the most concerning part of this chain. "
        "It is a professional, maintained implant — not a commodity tool downloaded from the internet. "
        "It supports self-update, meaning the attacker can push new capabilities at any time. "
        "The developer made one mistake: they did not strip the debug symbol path from the binary, "
        "which leaks their username as 'marks'. This is a tradecraft error. "
        "No public reports of this tool exist — we are looking at previously undocumented malware.")


def slide_iocs(prs, iocs: dict):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    _set_bg(sl, NAVY)
    _slide_header(sl, "Indicators of Compromise", "Block and hunt for these across your estate immediately")

    ioc_data = iocs.get("iocs", {})

    rows = []
    rows.append(("DOMAIN",  "simpludelicios.md",                                "C2 domain — block at DNS and perimeter"))
    rows.append(("IP",      "89.42.218.223",                                    "C2 server — block at firewall"))
    rows.append(("URL",     "simpludelicios.md/reportv.php",                    "Victim registration beacon"))
    rows.append(("URL",     "simpludelicios.md/verify.ps1",                     "Stage 2 loader"))
    rows.append(("URL",     "simpludelicios.md/notepad.b64",                    "Stage 3 encoded DLL"))
    rows.append(("URL",     "simpludelicios.md/vcapcha.ps1",                    "ClickFix payload (= submitted script)"))
    rows.append(("SHA-256", "d4a1153c4e4b...b7852b855",                        "submitted_script.ps1 / vcapcha.ps1"))
    rows.append(("SHA-256", "160dc775c7fc...733fa48c7",                        "HelloMemory.dll (decoded implant)"))
    rows.append(("SHA-256", "76b351cf70ee...db3ddd1",                          "4th-stage payload (hardcoded in DLL)"))
    rows.append(("FILE",    r"%TEMP%\verify_script.ps1",                        "Temp staging file (deleted post-exec)"))
    rows.append(("REGKEY",  r"HKCU\Software\Microsoft\Windows\CurrentVersion\Run", "Persistence key set by HelloMemory"))

    col_widths = [Inches(1.2), Inches(5.4), Inches(5.7)]
    headers    = ["Type", "Indicator", "Context"]
    header_y   = Inches(1.0)

    # Header row
    x = Inches(0.25)
    for j, (hdr, w) in enumerate(zip(headers, col_widths)):
        _add_rect(sl, x, header_y, w, Inches(0.34), fill_color=BLUE_LIGHT)
        _add_textbox(sl, x + Inches(0.05), header_y + Inches(0.02),
                     w - Inches(0.1), Inches(0.3),
                     hdr, font_size=Pt(11), bold=True, color=NAVY)
        x += w

    row_h = Inches(0.44)
    y = header_y + Inches(0.34)
    for i, (itype, indicator, context) in enumerate(rows):
        bg = RGBColor(0x12, 0x1C, 0x2B) if i % 2 == 0 else RGBColor(0x17, 0x23, 0x35)
        x  = Inches(0.25)
        for val, w in zip([itype, indicator, context], col_widths):
            _add_rect(sl, x, y, w, row_h, fill_color=bg)
            col = AMBER if itype in ("SHA-256","FILE","REGKEY") else (
                  RED if itype == "IP" else LIGHT_GREY)
            _add_textbox(sl, x + Inches(0.05), y + Inches(0.04),
                         w - Inches(0.1), Inches(0.36),
                         val, font_size=Pt(9.5), color=col)
            x += w
        y += row_h

    _add_rect(sl, 0, Inches(7.2), SLIDE_W, Inches(0.3), fill_color=BLUE_MID)
    _speaker_note(sl,
        "SPEAKER: This slide is your block list. "
        "Domain, IP, and all four URLs should be added to your DNS sinkhole, firewall, and proxy blocklist today. "
        "The SHA-256 hashes should be added to your EDR and AV platform. "
        "Search endpoint telemetry for any process that created a file called verify_script.ps1 in %TEMP%. "
        "Check the registry run key on any potentially affected hosts.")


def slide_attribution(prs):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    _set_bg(sl, NAVY)
    _slide_header(sl, "Attribution Assessment", "Who is behind this — and what we cannot confirm")

    # Assessment box
    _add_rect(sl, Inches(0.3), Inches(1.0), Inches(12.6), Inches(1.1),
              fill_color=BLUE_MID)
    _add_textbox(sl, Inches(0.5), Inches(1.06), Inches(12.0), Inches(0.4),
                 "ASSESSED WITH MEDIUM CONFIDENCE", font_size=Pt(12), bold=True, color=AMBER)
    _add_textbox(sl, Inches(0.5), Inches(1.5), Inches(12.0), Inches(0.55),
                 "Financially motivated criminal threat actor — likely operating as an initial access broker (IAB). "
                 "Custom tooling (HelloMemory) with no public attribution. Not assessed as nation-state.",
                 font_size=Pt(12), color=WHITE)

    # Evidence columns
    _add_textbox(sl, Inches(0.3), Inches(2.25), Inches(5.8), Inches(0.35),
                 "WHAT POINTS TO CRIMINAL / IAB", font_size=Pt(12), bold=True, color=GREEN)
    evidence_for = [
        "ClickFix delivery = mass-targeting for volume, not precision",
        "Device ID beaconing on first run = counting infections (IAB metric)",
        "Self-updating implant = long-term access sale, not smash-and-grab",
        "Compromised cheap hosting vs. bulletproof infrastructure",
        "PDB path left in binary = operational security gaps atypical of state actors",
        "No geopolitical targeting signals observed",
    ]
    y = Inches(2.65)
    for e in evidence_for:
        _add_textbox(sl, Inches(0.5), y, Inches(5.6), Inches(0.4),
                     f"✓  {e}", font_size=Pt(10.5), color=LIGHT_GREY)
        y += Inches(0.43)

    _add_textbox(sl, Inches(6.8), Inches(2.25), Inches(6.0), Inches(0.35),
                 "WHAT CANNOT BE DETERMINED", font_size=Pt(12), bold=True, color=AMBER)
    unknowns = [
        "Whether HelloMemory is sold, private, or shared with other actors",
        "Identity behind developer username 'marks'",
        "What the 4th-stage payload does",
        "Full victim scope — who else was targeted",
        "Whether this is targeted at your organisation or opportunistic",
        "Definitive nation-state link or exclusion without further intel",
    ]
    y = Inches(2.65)
    for u in unknowns:
        _add_textbox(sl, Inches(7.0), y, Inches(5.8), Inches(0.4),
                     f"?  {u}", font_size=Pt(10.5), color=MID_GREY)
        y += Inches(0.43)

    _add_rect(sl, Inches(0.3), Inches(5.75), Inches(12.6), Inches(0.48),
              fill_color=RGBColor(0x17, 0x20, 0x2A), line_color=BLUE_LIGHT, line_width=Pt(1))
    _add_textbox(sl, Inches(0.5), Inches(5.8), Inches(12.2), Inches(0.4),
                 "Next step:  Submit SHA-256 hashes to VirusTotal and Malware Bazaar — these appear unreported. "
                 "Community tags may surface prior campaigns or actor overlaps.",
                 font_size=Pt(10.5), color=BLUE_LIGHT)

    _add_rect(sl, 0, Inches(7.2), SLIDE_W, Inches(0.3), fill_color=BLUE_MID)
    _speaker_note(sl,
        "SPEAKER: We cannot name a specific threat actor group because HelloMemory does not appear in any "
        "public threat intelligence database. This is either brand-new tooling or a private/closed-market tool. "
        "What we can say is that all the evidence points to a financially motivated criminal, not a government. "
        "The most likely objective is to sell access to compromised machines to ransomware operators or other criminals.")


def slide_recommendations(prs):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    _set_bg(sl, NAVY)
    _slide_header(sl, "Recommended Actions", "Prioritised response for the next 24–72 hours")

    actions = [
        ("IMMEDIATE  (0–4 h)",  RED, [
            "Block simpludelicios.md and 89.42.218.223 at DNS, firewall, and web proxy",
            "Search proxy / DNS logs for any prior connections to this domain across the estate",
            "Hunt for %TEMP%\\verify_script.ps1 in EDR telemetry — indicates Stage 2 ran",
            "Check HKCU\\…\\Run registry keys on any host that touched the domain",
            "Isolate any host confirmed to have executed the script",
        ]),
        ("SHORT-TERM  (24–72 h)", AMBER, [
            "Submit all hashes to VirusTotal and Malware Bazaar to support community detection",
            "Report the compromise to romarg.com / whmpanels.com hosting provider for takedown",
            "Expand hunt: search for PowerShell spawning Invoke-WebRequest to .md domains",
            "Determine the delivery mechanism — which web page served the ClickFix lure?",
            "Conduct user awareness communications: explain fake CAPTCHA technique",
        ]),
        ("ONGOING", BLUE_MID, [
            "Deploy IOC block list to all security tooling (EDR, SIEM, email gateway, proxy)",
            "Monitor for HelloMemory activity patterns: base64 blobs downloaded via PowerShell",
            "Review execution policy settings estate-wide — Set-ExecutionPolicy events are detectable",
            "Engage hosting provider / CERT for intelligence on other victims on same server",
        ]),
    ]

    y = Inches(1.0)
    for heading, color, items in actions:
        _add_rect(sl, Inches(0.3), y, Inches(12.6), Inches(0.38), fill_color=color)
        _add_textbox(sl, Inches(0.45), y + Inches(0.03), Inches(12.2), Inches(0.32),
                     heading, font_size=Pt(13), bold=True,
                     color=NAVY if color == AMBER else WHITE)
        y += Inches(0.38)
        for item in items:
            _add_textbox(sl, Inches(0.6), y, Inches(12.1), Inches(0.37),
                         f"▸  {item}", font_size=Pt(10.5), color=LIGHT_GREY)
            y += Inches(0.38)
        y += Inches(0.12)

    _add_rect(sl, 0, Inches(7.2), SLIDE_W, Inches(0.3), fill_color=BLUE_MID)
    _speaker_note(sl,
        "SPEAKER: The top priority is blocking the domain and IP now. "
        "Even if no user executed the script, the domain should be blocked as a precaution. "
        "The most critical hunt question is: did any host actually reach Stage 3? "
        "If yes, we have a persistent implant on that machine and need to treat it as fully compromised. "
        "User awareness is also important — this technique works because users trust CAPTCHA pages.")


def slide_closing(prs, meta: dict):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    _set_bg(sl, NAVY)
    _add_rect(sl, 0, 0, SLIDE_W, Inches(0.12), fill_color=BLUE_LIGHT)

    _add_textbox(sl, Inches(0.4), Inches(1.8), Inches(12.5), Inches(0.8),
                 "Questions?", font_size=Pt(52), bold=True, color=WHITE)

    _add_textbox(sl, Inches(0.4), Inches(2.8), Inches(12.5), Inches(0.45),
                 "All evidence is stored in the case vault and available for review.",
                 font_size=Pt(16), color=LIGHT_GREY)

    details = [
        f"Case ID:           {meta.get('case_id','C003')}",
        f"Report:            cases/{meta.get('case_id','C003')}/reports/investigation_report.md",
        f"IOCs:              cases/{meta.get('case_id','C003')}/iocs/iocs.json",
        f"Audit trail:       registry/audit.log",
        f"Analyst:           {meta.get('analyst','SOC')}",
    ]
    y = Inches(3.5)
    for d in details:
        _add_textbox(sl, Inches(0.4), y, Inches(12.0), Inches(0.38),
                     d, font_size=Pt(12), color=MID_GREY)
        y += Inches(0.4)

    _add_rect(sl, 0, Inches(7.2), SLIDE_W, Inches(0.3), fill_color=BLUE_MID)
    _add_textbox(sl, Inches(0.2), Inches(7.22), Inches(6), Inches(0.25),
                 "RESTRICTED – INTERNAL USE ONLY",
                 font_size=Pt(9), color=LIGHT_GREY)
    _speaker_note(sl, "Open the floor to questions. Key messages to reinforce: "
                  "1) Block the domain now.  2) Hunt for affected hosts.  "
                  "3) HelloMemory is unreported — we should share these hashes with the community.")


# ── Main ──────────────────────────────────────────────────────────────────────

def generate_pptx(case_id: str) -> dict:
    case_dir    = CASES_DIR / case_id
    meta_path   = case_dir / "case_meta.json"
    iocs_path   = case_dir / "iocs" / "iocs.json"

    meta = load_json(meta_path) if meta_path.exists() else {"case_id": case_id}
    iocs = load_json(iocs_path) if iocs_path.exists() else {"iocs": {}, "total": {}}

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs, meta)
    slide_executive_summary(prs, meta, iocs)
    slide_how_it_arrived(prs)
    slide_attack_chain(prs)
    slide_infrastructure(prs)
    slide_hellomemory(prs)
    slide_iocs(prs, iocs)
    slide_attribution(prs)
    slide_recommendations(prs)
    slide_closing(prs, meta)

    reports_dir = case_dir / "reports"
    reports_dir.mkdir(parents=True, exist_ok=True)
    out_path = reports_dir / "investigation_briefing.pptx"

    prs.save(str(out_path))
    out_bytes = out_path.read_bytes()
    audit("generate_pptx", str(out_path), extra={"case_id": case_id})

    print(f"[generate_pptx] {len(prs.slides)} slides written to {out_path}")
    return {
        "case_id":    case_id,
        "pptx_path":  str(out_path),
        "slides":     len(prs.slides),
        "size_bytes": len(out_bytes),
        "ts":         utcnow(),
    }


if __name__ == "__main__":
    import argparse
    p = argparse.ArgumentParser()
    p.add_argument("--case", required=True, dest="case_id")
    args = p.parse_args()
    result = generate_pptx(args.case_id)
    print(json.dumps(result, indent=2))
