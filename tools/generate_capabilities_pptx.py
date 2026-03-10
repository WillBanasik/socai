"""
tool: generate_capabilities_pptx
---------------------------------
Generates a platform capabilities showcase PowerPoint deck using a real
case (IV_CASE_003 — ClickFix / HelloMemory) as the running example.

Demonstrates the full SOCAI pipeline: IOC collection, enrichment,
verdict scoring, hunt query generation, and YAML export.

Usage:
  python3 tools/generate_capabilities_pptx.py

Output:
  socai_capabilities_showcase.pptx  (in repo root)
"""
from __future__ import annotations

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x0D, 0x1B, 0x2A)
DARK_NAVY   = RGBColor(0x09, 0x12, 0x1C)
BLUE_MID    = RGBColor(0x1B, 0x4F, 0x72)
BLUE_LIGHT  = RGBColor(0x21, 0x8D, 0xBF)
BLUE_ACCENT = RGBColor(0x3A, 0xA6, 0xD5)
AMBER       = RGBColor(0xF3, 0x9C, 0x12)
RED         = RGBColor(0xC0, 0x39, 0x2B)
GREEN       = RGBColor(0x1E, 0x8B, 0x4C)
PURPLE      = RGBColor(0x6E, 0x27, 0x94)
TEAL        = RGBColor(0x17, 0xA5, 0x89)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY  = RGBColor(0xD5, 0xD8, 0xDC)
MID_GREY    = RGBColor(0x7F, 0x8C, 0x8D)
DARK_BG     = RGBColor(0x12, 0x1C, 0x2B)
DARKER_BG   = RGBColor(0x17, 0x20, 0x2A)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Low-level helpers ─────────────────────────────────────────────────────────

def _set_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # RECTANGLE
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def _add_textbox(slide, l, t, w, h, text, font_size=Pt(14),
                 color=WHITE, bold=False, align=PP_ALIGN.LEFT,
                 italic=False, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = font_size
    run.font.color.rgb = color
    run.font.bold  = bold
    run.font.italic = italic
    return txb


def _add_multiline(slide, l, t, w, h, lines, font_size=Pt(11),
                   color=LIGHT_GREY, bold=False, spacing=Inches(0)):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size = font_size
        run.font.color.rgb = color
        run.font.bold = bold
    return txb


def _add_label_box(slide, l, t, w, h, label, value,
                   box_color=BLUE_MID, label_size=Pt(10), val_size=Pt(16)):
    rect = _add_rect(slide, l, t, w, h, fill_color=box_color)
    tf   = rect.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = label.upper()
    r1.font.size  = label_size
    r1.font.color.rgb = LIGHT_GREY
    r1.font.bold  = False
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = value
    r2.font.size  = val_size
    r2.font.color.rgb = WHITE
    r2.font.bold  = True
    return rect


def _add_flow_box(slide, l, t, w, h, title, lines,
                  box_fill=BLUE_MID, title_size=Pt(13), body_size=Pt(10)):
    rect = _add_rect(slide, l, t, w, h, fill_color=box_fill)
    tf   = rect.text_frame
    tf.word_wrap = True
    tf.margin_left   = Inches(0.12)
    tf.margin_right  = Inches(0.12)
    tf.margin_top    = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    r0 = p0.add_run()
    r0.text = title
    r0.font.size  = title_size
    r0.font.color.rgb = WHITE
    r0.font.bold  = True
    for line in lines:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line
        r.font.size  = body_size
        r.font.color.rgb = LIGHT_GREY
    return rect


def _slide_header(slide, title: str, subtitle: str = ""):
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.85), fill_color=BLUE_MID)
    _add_textbox(slide, Inches(0.3), Inches(0.08), Inches(10), Inches(0.55),
                 title, font_size=Pt(24), bold=True, color=WHITE)
    if subtitle:
        _add_textbox(slide, Inches(0.3), Inches(0.6), Inches(10), Inches(0.3),
                     subtitle, font_size=Pt(12), color=LIGHT_GREY)


def _footer(slide):
    _add_rect(slide, 0, Inches(7.2), SLIDE_W, Inches(0.3), fill_color=BLUE_MID)
    _add_textbox(slide, Inches(0.2), Inches(7.22), Inches(6), Inches(0.25),
                 "SOCAI  |  Automated SOC Investigation Platform",
                 font_size=Pt(9), color=LIGHT_GREY)


def _speaker_note(slide, text: str):
    notes = slide.notes_slide
    tf    = notes.notes_text_frame
    tf.text = text


def _code_block(slide, l, t, w, h, text, font_size=Pt(9), color=BLUE_LIGHT):
    """Simulated monospace code block."""
    rect = _add_rect(slide, l, t, w, h,
                     fill_color=RGBColor(0x0A, 0x14, 0x1E),
                     line_color=BLUE_MID, line_width=Pt(0.75))
    _add_textbox(slide, l + Inches(0.1), t + Inches(0.06),
                 w - Inches(0.2), h - Inches(0.12),
                 text, font_size=font_size, color=color)
    return rect


# ── Slide 1: Title ───────────────────────────────────────────────────────────

def slide_title(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, NAVY)
    _add_rect(sl, 0, 0, SLIDE_W, Inches(0.12), fill_color=BLUE_LIGHT)

    _add_textbox(sl, Inches(0.5), Inches(1.4), Inches(12.5), Inches(1.2),
                 "SOCAI", font_size=Pt(64), bold=True, color=WHITE)

    _add_textbox(sl, Inches(0.5), Inches(2.6), Inches(12.5), Inches(0.7),
                 "Automated SOC Investigation Platform",
                 font_size=Pt(28), color=BLUE_LIGHT)

    _add_textbox(sl, Inches(0.5), Inches(3.5), Inches(12.5), Inches(0.5),
                 "From alert to hunt queries in minutes  —  not hours",
                 font_size=Pt(16), color=LIGHT_GREY, italic=True)

    # Feature pills
    pills = [
        "IOC Extraction", "Multi-Provider Enrichment", "Verdict Scoring",
        "Hunt Query Generation", "YAML Export", "Campaign Clustering",
    ]
    pill_w = Inches(2.0)
    pill_h = Inches(0.42)
    pill_t = Inches(4.4)
    for i, pill in enumerate(pills):
        col_idx = i % 3
        row_idx = i // 3
        pl = Inches(0.5) + col_idx * (pill_w + Inches(0.2))
        pt_y = pill_t + row_idx * (pill_h + Inches(0.15))
        _add_rect(sl, pl, pt_y, pill_w, pill_h,
                  fill_color=BLUE_MID, line_color=BLUE_LIGHT, line_width=Pt(1))
        _add_textbox(sl, pl, pt_y + Inches(0.04), pill_w, pill_h - Inches(0.08),
                     pill, font_size=Pt(11), color=WHITE, bold=True,
                     align=PP_ALIGN.CENTER)

    # Case study callout
    _add_rect(sl, Inches(7.0), Inches(4.3), Inches(5.8), Inches(1.15),
              fill_color=DARKER_BG, line_color=RED, line_width=Pt(1.5))
    _add_textbox(sl, Inches(7.2), Inches(4.35), Inches(5.4), Inches(0.35),
                 "LIVE CASE STUDY", font_size=Pt(12), bold=True, color=RED)
    _add_textbox(sl, Inches(7.2), Inches(4.7), Inches(5.4), Inches(0.7),
                 "IV_CASE_003 — ClickFix / HelloMemory\n"
                 "Multi-stage malware delivery via fake CAPTCHA\n"
                 "Threat actor: Financially motivated IAB",
                 font_size=Pt(11), color=LIGHT_GREY)

    _footer(sl)
    _speaker_note(sl,
        "SOCAI is an automated security operations investigation platform. "
        "Today we walk through a real case — IV_CASE_003, a ClickFix social engineering "
        "attack delivering the HelloMemory implant — to demonstrate how the platform "
        "collects IOCs, enriches them, scores verdicts, generates hunt queries, "
        "and exports everything as structured YAML.")


# ── Slide 2: Pipeline Overview ───────────────────────────────────────────────

def slide_pipeline(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, NAVY)
    _slide_header(sl, "Investigation Pipeline", "15-step orchestrated flow from alert to actionable intelligence")

    # Pipeline stages in 3 rows
    stages = [
        # Row 1: Intake + Parallel
        [
            ("1. Case Create", BLUE_MID, ["Metadata", "Severity", "Registry"]),
            ("2. Triage", TEAL, ["IOC index check", "Known-malicious?", "Severity escalation"]),
            ("3. Planner", BLUE_MID, ["Route selection", "Input analysis"]),
            ("4. Email Analyst", PURPLE, ["Header parse", "SPF/DKIM/DMARC", "URL extraction"]),
        ],
        # Row 2: Parallel agents + Enrichment
        [
            ("5. PARALLEL", RED, ["Domain Investigator", "File Analyst", "Log Correlator"]),
            ("6. Sandbox", AMBER, ["Hash lookup", "Live detonation", "TTP extraction"]),
            ("7. Recursive Capture", BLUE_MID, ["Depth 2-N crawl", "New URL discovery"]),
            ("8. Phishing Detect", PURPLE, ["Brand impersonation", "LLM vision scan"]),
        ],
        # Row 3: Enrichment + Output
        [
            ("9. Enrichment", TEAL, ["15+ providers", "Parallel execution", "Cache + TTL"]),
            ("10. Verdict Score", GREEN, ["Composite scoring", "Confidence levels"]),
            ("11. Campaign", BLUE_MID, ["Cross-case IOC", "Union-Find cluster"]),
            ("12. Report + Queries", AMBER, ["Markdown report", "KQL/SPL/LogScale", "YAML export"]),
        ],
    ]

    bw = Inches(2.85)
    bh = Inches(1.5)
    gap_x = Inches(0.28)
    gap_y = Inches(0.18)
    start_y = Inches(1.0)

    for row_idx, row in enumerate(stages):
        for col_idx, (title, color, lines) in enumerate(row):
            bl = Inches(0.25) + col_idx * (bw + gap_x)
            bt = start_y + row_idx * (bh + gap_y)
            _add_flow_box(sl, bl, bt, bw, bh, title, lines,
                          box_fill=color, title_size=Pt(11), body_size=Pt(9))

    # Arrow annotations between rows
    for row_idx in range(2):
        y = start_y + (row_idx + 1) * (bh + gap_y) - gap_y + Inches(0.01)
        _add_textbox(sl, Inches(5.5), y - Inches(0.08), Inches(2), Inches(0.16),
                     "v  v  v  v", font_size=Pt(10), color=BLUE_LIGHT,
                     align=PP_ALIGN.CENTER)

    # Callout
    _add_rect(sl, Inches(0.25), Inches(6.1), Inches(12.7), Inches(0.65),
              fill_color=DARKER_BG, line_color=BLUE_LIGHT, line_width=Pt(1))
    _add_textbox(sl, Inches(0.45), Inches(6.15), Inches(12.3), Inches(0.55),
                 "Step 5 runs Domain/File/Log agents in parallel via ThreadPoolExecutor(max_workers=3).  "
                 "Enrichment (step 9) runs all provider calls in parallel (default 10 workers).  "
                 "Each step is fault-isolated — a failing step does not abort subsequent steps.",
                 font_size=Pt(10), color=LIGHT_GREY)

    _footer(sl)
    _speaker_note(sl,
        "The pipeline runs 15 steps from case creation to final report. "
        "Key architectural decisions: parallel execution for independent agents, "
        "fault isolation per step, cross-case intelligence via the IOC index, "
        "and idempotent re-runs so you can safely re-investigate.")


# ── Slide 3: Threat Actor / Case Study ───────────────────────────────────────

def slide_threat_actor(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, NAVY)
    _slide_header(sl, "Case Study: ClickFix / HelloMemory",
                  "IV_CASE_003 — Multi-stage malware delivery via fake CAPTCHA social engineering")

    # Attack chain (horizontal flow)
    chain = [
        ("STAGE 0\nCLICKFIX LURE", PURPLE,
         ["Fake CAPTCHA page", "Clipboard hijack", "Win+R  Ctrl+V  Enter"]),
        ("STAGE 1\nDROPPER", AMBER,
         ["submitted_script.ps1", "Beacon to C2", "Download Stage 2"]),
        ("STAGE 2\nLOADER", RGBColor(0xCA, 0x6F, 0x1E),
         ["verify.ps1", "Base64 decode in RAM", "Reflective .NET load"]),
        ("STAGE 3\nIMPLANT", RED,
         ["HelloMemory.dll", "Registry persistence", "Self-update + C2"]),
    ]

    bw = Inches(2.75)
    bh = Inches(2.6)
    bt = Inches(1.0)
    gap = Inches(0.3)
    for i, (title, color, lines) in enumerate(chain):
        bl = Inches(0.3) + i * (bw + gap)
        _add_flow_box(sl, bl, bt, bw, bh, title, lines,
                      box_fill=color, title_size=Pt(12), body_size=Pt(10.5))
        if i < 3:
            ax = bl + bw + Inches(0.02)
            ay = bt + bh / 2 - Inches(0.15)
            _add_textbox(sl, ax, ay, gap - Inches(0.04), Inches(0.3),
                         ">", font_size=Pt(26), color=BLUE_LIGHT,
                         align=PP_ALIGN.CENTER)

    # Threat actor profile
    _add_rect(sl, Inches(0.3), Inches(3.85), Inches(6.3), Inches(2.4),
              fill_color=DARK_BG, line_color=BLUE_MID, line_width=Pt(1))
    _add_textbox(sl, Inches(0.5), Inches(3.9), Inches(5.9), Inches(0.35),
                 "THREAT ACTOR PROFILE", font_size=Pt(12), bold=True, color=BLUE_LIGHT)
    profile_lines = [
        "Assessment:     Financially motivated criminal (IAB)",
        "Confidence:     MEDIUM",
        "Infrastructure: Compromised Romanian shared hosting",
        "                (89.42.218.223 / romarg srl / 6 live CVEs)",
        "Tooling:        Custom HelloMemory .NET implant (unreported)",
        "Developer:      Username 'marks' leaked via PDB path",
        "Objective:      Mass initial access for resale",
    ]
    y = Inches(4.3)
    for line in profile_lines:
        _add_textbox(sl, Inches(0.5), y, Inches(5.9), Inches(0.28),
                     line, font_size=Pt(9.5), color=LIGHT_GREY)
        y += Inches(0.28)

    # Key stats
    stats = [
        ("IOCs Extracted", "20"),
        ("C2 Endpoints", "4"),
        ("File Hashes", "5"),
        ("CVEs on Server", "6"),
    ]
    sx = Inches(6.85)
    sy = Inches(3.85)
    sw = Inches(3.1)
    sh = Inches(0.92)
    for i, (label, val) in enumerate(stats):
        _add_label_box(sl, sx, sy + i * (sh + Inches(0.12)), sw, sh,
                       label, val, box_color=BLUE_MID,
                       label_size=Pt(9), val_size=Pt(20))

    _footer(sl)
    _speaker_note(sl,
        "IV_CASE_003 was triggered by submission of a suspicious PowerShell script. "
        "The investigation revealed a four-stage attack chain originating from "
        "a ClickFix fake CAPTCHA campaign. The final payload, HelloMemory.dll, "
        "is a custom .NET implant that had no public threat intelligence matches — "
        "this is unreported tooling. The threat actor is assessed as a criminal "
        "initial access broker operating from compromised Romanian infrastructure.")


# ── Slide 4: IOC Collection ──────────────────────────────────────────────────

def slide_ioc_collection(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, NAVY)
    _slide_header(sl, "IOC Collection Engine",
                  "Automated extraction from 7+ artefact sources across all investigation inputs")

    # Sources panel (left)
    sources = [
        ("Web Captures",      "page.html, page.txt, redirect chains, XHR responses"),
        ("Static File Analysis", "PE headers, strings, entropy, embedded URLs"),
        ("Email Parsing",     ".eml headers, body URLs, attachment hashes"),
        ("Log Correlation",   "Parsed CSV/JSON logs, entity extraction"),
        ("Sandbox Results",   "Network IOCs, C2 beacons from detonation"),
        ("Recursive Crawl",   "Depth 2-N URL discovery from extracted links"),
    ]

    _add_textbox(sl, Inches(0.3), Inches(0.95), Inches(6.0), Inches(0.35),
                 "EXTRACTION SOURCES", font_size=Pt(13), bold=True, color=BLUE_LIGHT)
    y = Inches(1.35)
    for name, desc in sources:
        _add_rect(sl, Inches(0.3), y, Inches(2.3), Inches(0.42), fill_color=BLUE_MID)
        _add_textbox(sl, Inches(0.4), y + Inches(0.05), Inches(2.1), Inches(0.32),
                     name, font_size=Pt(10), bold=True, color=WHITE)
        _add_textbox(sl, Inches(2.75), y + Inches(0.05), Inches(4.3), Inches(0.32),
                     desc, font_size=Pt(9.5), color=LIGHT_GREY)
        y += Inches(0.5)

    # IOC types panel (right)
    _add_textbox(sl, Inches(7.3), Inches(0.95), Inches(5.7), Inches(0.35),
                 "IOC TYPES EXTRACTED", font_size=Pt(13), bold=True, color=BLUE_LIGHT)

    ioc_types = [
        ("IPv4", "1", "89.42.218.223"),
        ("Domain", "3", "simpludelicios.md"),
        ("URL", "5", "simpludelicios.md/reportv.php"),
        ("MD5", "3", "f41a63eb3f7b2794..."),
        ("SHA1", "3", "c0b3f926a239bd30..."),
        ("SHA256", "5", "160dc775c7fccd4a..."),
        ("Email", "0", "--"),
        ("CVE", "0", "--"),
    ]

    col_widths = [Inches(0.9), Inches(0.7), Inches(3.9)]
    hdr_y = Inches(1.35)
    x0 = Inches(7.3)

    # Header
    for j, (hdr, w) in enumerate(zip(["Type", "Count", "Example (IV_CASE_003)"], col_widths)):
        _add_rect(sl, x0 + sum(cw for cw in [Inches(0)] + list(col_widths[:j])),
                  hdr_y, w, Inches(0.32), fill_color=BLUE_LIGHT)
        _add_textbox(sl, x0 + sum(cw for cw in [Inches(0)] + list(col_widths[:j])) + Inches(0.05),
                     hdr_y + Inches(0.02), w - Inches(0.1), Inches(0.28),
                     hdr, font_size=Pt(9), bold=True, color=NAVY)

    y = hdr_y + Inches(0.32)
    for i, (itype, count, example) in enumerate(ioc_types):
        bg = DARK_BG if i % 2 == 0 else DARKER_BG
        cx = x0
        for val, w in zip([itype, count, example], col_widths):
            _add_rect(sl, cx, y, w, Inches(0.34), fill_color=bg)
            col = AMBER if int(count) > 0 and val == count else LIGHT_GREY
            _add_textbox(sl, cx + Inches(0.05), y + Inches(0.03), w - Inches(0.1),
                         Inches(0.28), val, font_size=Pt(9), color=col)
            cx += w
        y += Inches(0.34)

    # Regex + de-duplication callout
    _add_rect(sl, Inches(0.3), Inches(4.65), Inches(12.7), Inches(1.4),
              fill_color=DARK_BG, line_color=BLUE_MID, line_width=Pt(1))
    _add_textbox(sl, Inches(0.5), Inches(4.7), Inches(12.2), Inches(0.35),
                 "HOW IT WORKS", font_size=Pt(12), bold=True, color=BLUE_LIGHT)

    code = (
        "python3 socai.py investigate --case IV_CASE_003 --title \"Suspicious PowerShell\" "
        "--severity high --zip sample.zip --zip-pass infected --url \"https://simpludelicios.md/vcapcha.ps1\""
    )
    _code_block(sl, Inches(0.5), Inches(5.1), Inches(12.2), Inches(0.5),
                code, font_size=Pt(9))

    _add_textbox(sl, Inches(0.5), Inches(5.7), Inches(12.2), Inches(0.3),
                 "extract_iocs scans artefacts/ and logs/ directories. "
                 "IOCs are de-duplicated, typed, and written to iocs/iocs.json. "
                 "7 artefact files scanned in IV_CASE_003, yielding 20 unique IOCs across 6 types.",
                 font_size=Pt(10), color=MID_GREY)

    _footer(sl)
    _speaker_note(sl,
        "IOC extraction is fully automated and source-agnostic. "
        "It uses compiled regex patterns for each IOC type and scans every "
        "artefact file produced by the investigation. The output is a canonical "
        "JSON file that feeds the enrichment and query generation stages.")


# ── Slide 5: Enrichment ─────────────────────────────────────────────────────

def slide_enrichment(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, NAVY)
    _slide_header(sl, "Multi-Provider Enrichment",
                  "Parallel intelligence lookups across 15+ threat intelligence providers")

    # Provider grid
    providers = [
        ("VirusTotal",      "IP, domain, URL, hash",  BLUE_MID),
        ("AbuseIPDB",       "IPv4",                    BLUE_MID),
        ("Shodan",          "IPv4",                    BLUE_MID),
        ("GreyNoise",       "IPv4",                    BLUE_MID),
        ("URLScan.io",      "Domain, URL",             BLUE_MID),
        ("Intezer",         "MD5, SHA1, SHA256",       BLUE_MID),
        ("OpenCTI",         "All types + CVE",         PURPLE),
        ("AlienVault OTX",  "IP, domain, URL, hash",   PURPLE),
        ("Hybrid Analysis", "SHA256",                  PURPLE),
        ("Censys",          "IPv4, domain",            PURPLE),
        ("WHOISXML",        "Domain (age, registrant)", TEAL),
        ("proxycheck.io",   "IPv4",                    TEAL),
        ("EmailRep.io",     "Email",                   TEAL),
        ("Abuse.ch",        "IP, domain, URL, hash",   TEAL),
        ("Any.Run",         "SHA256 (sandbox)",        RED),
        ("Joe Sandbox",     "SHA256 (sandbox)",        RED),
    ]

    cols = 4
    pw = Inches(2.95)
    ph = Inches(0.65)
    pgx = Inches(0.2)
    pgy = Inches(0.12)
    start_x = Inches(0.25)
    start_y = Inches(1.0)

    for i, (name, types, color) in enumerate(providers):
        col = i % cols
        row = i // cols
        px = start_x + col * (pw + pgx)
        py = start_y + row * (ph + pgy)
        _add_rect(sl, px, py, pw, ph, fill_color=color)
        _add_textbox(sl, px + Inches(0.1), py + Inches(0.04), pw - Inches(0.2),
                     Inches(0.3), name, font_size=Pt(11), bold=True, color=WHITE)
        _add_textbox(sl, px + Inches(0.1), py + Inches(0.32), pw - Inches(0.2),
                     Inches(0.25), types, font_size=Pt(8.5), color=LIGHT_GREY)

    # Architecture callout
    features = [
        ("Parallel Execution",  "ThreadPoolExecutor with configurable workers (default 10)"),
        ("Cross-Run Caching",   "registry/enrichment_cache.json with configurable TTL (default 24h)"),
        ("Skip-Enrichment",     "Triage pre-checks IOC index; skips already-covered IOCs"),
        ("Fault Isolation",     "Provider failures don't block other lookups"),
    ]

    y = Inches(4.2)
    _add_textbox(sl, Inches(0.3), y, Inches(12.5), Inches(0.35),
                 "ENRICHMENT ARCHITECTURE", font_size=Pt(13), bold=True, color=BLUE_LIGHT)
    y += Inches(0.4)
    for feat, desc in features:
        _add_rect(sl, Inches(0.3), y, Inches(2.6), Inches(0.38), fill_color=TEAL)
        _add_textbox(sl, Inches(0.4), y + Inches(0.04), Inches(2.4), Inches(0.3),
                     feat, font_size=Pt(10), bold=True, color=WHITE)
        _add_textbox(sl, Inches(3.1), y + Inches(0.04), Inches(9.5), Inches(0.3),
                     desc, font_size=Pt(10), color=LIGHT_GREY)
        y += Inches(0.45)

    _footer(sl)
    _speaker_note(sl,
        "Enrichment runs all provider calls in parallel. Results are cached "
        "cross-run to avoid redundant API calls. The triage stage pre-checks "
        "the IOC index for already-enriched IOCs and skips them. Each provider "
        "returns a normalised result with status, verdict, and provider-specific fields.")


# ── Slide 6: Verdict Scoring ────────────────────────────────────────────────

def slide_verdict_scoring(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, NAVY)
    _slide_header(sl, "Verdict Scoring Engine",
                  "Composite multi-provider verdict with confidence classification")

    # Scoring logic
    _add_textbox(sl, Inches(0.3), Inches(1.0), Inches(6.0), Inches(0.35),
                 "SCORING LOGIC", font_size=Pt(13), bold=True, color=BLUE_LIGHT)

    rules = [
        ("MALICIOUS",  RED,    "1+ provider says malicious AND malicious >= suspicious"),
        ("SUSPICIOUS", AMBER,  "1+ provider says suspicious AND malicious == 0"),
        ("CLEAN",      GREEN,  "All responsive providers say clean"),
    ]
    y = Inches(1.45)
    for label, color, desc in rules:
        _add_rect(sl, Inches(0.3), y, Inches(1.8), Inches(0.42), fill_color=color)
        _add_textbox(sl, Inches(0.4), y + Inches(0.06), Inches(1.6), Inches(0.3),
                     label, font_size=Pt(11), bold=True,
                     color=NAVY if color == AMBER else WHITE, align=PP_ALIGN.CENTER)
        _add_textbox(sl, Inches(2.3), y + Inches(0.06), Inches(4.5), Inches(0.3),
                     desc, font_size=Pt(10.5), color=LIGHT_GREY)
        y += Inches(0.52)

    # Confidence levels
    _add_textbox(sl, Inches(0.3), Inches(3.1), Inches(6.0), Inches(0.35),
                 "CONFIDENCE LEVELS", font_size=Pt(13), bold=True, color=BLUE_LIGHT)
    conf_rules = [
        ("HIGH",   TEAL,       "3+ providers, >66% agree"),
        ("MEDIUM", BLUE_MID,   "2+ providers, strict majority >50%"),
        ("LOW",    MID_GREY,   "1 provider, or 50/50 split"),
    ]
    y = Inches(3.5)
    for label, color, desc in conf_rules:
        _add_rect(sl, Inches(0.3), y, Inches(1.3), Inches(0.38), fill_color=color)
        _add_textbox(sl, Inches(0.4), y + Inches(0.04), Inches(1.1), Inches(0.3),
                     label, font_size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _add_textbox(sl, Inches(1.8), y + Inches(0.04), Inches(3.8), Inches(0.3),
                     desc, font_size=Pt(10), color=LIGHT_GREY)
        y += Inches(0.46)

    # Right side: Example verdict output
    _add_textbox(sl, Inches(7.0), Inches(1.0), Inches(6.0), Inches(0.35),
                 "EXAMPLE: VERDICT SUMMARY OUTPUT", font_size=Pt(13), bold=True, color=BLUE_LIGHT)

    verdict_yaml = (
        'verdict_summary.json:\n'
        '  ioc_count: 17\n'
        '  high_priority: []        # malicious IOCs\n'
        '  needs_review:\n'
        '    - "149.72.92.159"      # suspicious\n'
        '  clean:\n'
        '    - "www.googletagmanager.com"\n'
        '    - "cportal.hallswater.com"\n'
        '    - "www.webflex.biz"\n'
        '    - ... (14 more)\n'
        '\n'
        '  iocs:\n'
        '    "149.72.92.159":\n'
        '      verdict: suspicious\n'
        '      confidence: LOW\n'
        '      providers:\n'
        '        abuseipdb: clean\n'
        '        proxycheck: clean\n'
        '        virustotal: suspicious  # 1/3 = LOW'
    )
    _code_block(sl, Inches(7.0), Inches(1.4), Inches(5.9), Inches(3.7),
                verdict_yaml, font_size=Pt(9))

    # Auto-disposition callout
    _add_rect(sl, Inches(0.3), Inches(5.2), Inches(12.7), Inches(0.85),
              fill_color=DARKER_BG, line_color=GREEN, line_width=Pt(1.5))
    _add_textbox(sl, Inches(0.5), Inches(5.25), Inches(12.2), Inches(0.35),
                 "AUTO-DISPOSITION", font_size=Pt(12), bold=True, color=GREEN)
    _add_textbox(sl, Inches(0.5), Inches(5.6), Inches(12.2), Inches(0.4),
                 "After enrichment, if verdict_summary has 0 malicious and 0 suspicious IOCs, "
                 "the case is auto-closed with disposition 'benign_auto_closed' — unless the "
                 "report confidence score meets or exceeds SOCAI_CONF_AUTO_CLOSE (default 0.20).",
                 font_size=Pt(10), color=LIGHT_GREY)

    # IOC index callout
    _add_rect(sl, Inches(0.3), Inches(6.15), Inches(12.7), Inches(0.6),
              fill_color=DARKER_BG, line_color=AMBER, line_width=Pt(1))
    _add_textbox(sl, Inches(0.5), Inches(6.2), Inches(12.2), Inches(0.5),
                 "Cross-case IOC index (registry/ioc_index.json) tracks first_seen, last_seen, "
                 "cases[], and composite verdict per IOC. Recurring infrastructure triggers warnings.",
                 font_size=Pt(10), color=AMBER)

    _footer(sl)
    _speaker_note(sl,
        "The verdict scoring engine aggregates results from all providers. "
        "It requires consensus — a single suspicious flag doesn't override "
        "multiple clean results unless it reaches the malicious threshold. "
        "The IOC index provides cross-case memory so recurring infrastructure "
        "is flagged immediately on the next investigation.")


# ── Slide 7: Hunt Query Generation ──────────────────────────────────────────

def slide_hunt_queries(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, NAVY)
    _slide_header(sl, "Hunt Query Generation",
                  "Automated SIEM queries across KQL, Splunk SPL, and CrowdStrike LogScale")

    # Three platform columns
    kql_example = (
        'let suspect_ips = dynamic(\n'
        '  ["89.42.218.223"]);\n'
        'DeviceNetworkEvents\n'
        '| where RemoteIP in (suspect_ips)\n'
        '    or LocalIP in (suspect_ips)\n'
        '| project TimeGenerated,\n'
        '    DeviceName,\n'
        '    InitiatingProcessFileName,\n'
        '    RemoteIP, RemotePort,\n'
        '    RemoteUrl, ActionType\n'
        '| order by TimeGenerated asc'
    )

    splunk_example = (
        'index=* (src_ip IN (\n'
        '  89.42.218.223)\n'
        '  OR dest_ip IN (\n'
        '  89.42.218.223))\n'
        '| eval direction=if(\n'
        '    src_ip IN (\n'
        '    89.42.218.223),\n'
        '    "outbound","inbound")\n'
        '| table _time, src_ip,\n'
        '    src_host, dest_ip,\n'
        '    dest_port, direction\n'
        '| sort _time'
    )

    logscale_example = (
        '#event_simpleName =\n'
        '  NetworkConnectIP4\n'
        '| RemoteAddressIP4 =\n'
        '    "89.42.218.223"\n'
        '| table(\n'
        '    [@timestamp,\n'
        '     ComputerName,\n'
        '     LocalAddressIP4,\n'
        '     RemoteAddressIP4,\n'
        '     RemotePort,\n'
        '     FileName,\n'
        '     CommandLine])'
    )

    platforms = [
        ("KQL — Sentinel / MDE",  kql_example,    BLUE_LIGHT),
        ("Splunk SPL",            splunk_example,  AMBER),
        ("LogScale / Falcon",     logscale_example, GREEN),
    ]

    pw = Inches(3.9)
    gap = Inches(0.25)
    for i, (title, code, accent) in enumerate(platforms):
        px = Inches(0.25) + i * (pw + gap)

        _add_rect(sl, px, Inches(1.0), pw, Inches(0.4), fill_color=accent)
        _add_textbox(sl, px + Inches(0.1), Inches(1.04), pw - Inches(0.2),
                     Inches(0.32), title, font_size=Pt(12), bold=True,
                     color=NAVY if accent == AMBER else WHITE, align=PP_ALIGN.CENTER)

        _code_block(sl, px, Inches(1.48), pw, Inches(3.1), code, font_size=Pt(8.5))

    # Query types
    _add_textbox(sl, Inches(0.3), Inches(4.75), Inches(12.5), Inches(0.35),
                 "QUERY CATEGORIES GENERATED", font_size=Pt(13), bold=True, color=BLUE_LIGHT)

    categories = [
        ("IPv4 Lookups",    "DeviceNetworkEvents, IdentityLogonEvents, SecurityEvent, Syslog, BehaviorAnalytics"),
        ("Domain Lookups",  "DeviceNetworkEvents (RemoteUrl), IdentityQueryEvents"),
        ("Hash Lookups",    "DeviceFileEvents, DeviceProcessEvents, DeviceImageLoadEvents, AlertEvidence"),
        ("URL Lookups",     "DeviceNetworkEvents (RemoteUrl has)"),
        ("Email Lookups",   "EmailEvents, EmailUrlInfo"),
        ("Unified Timeline","Union across all tables — single chronological view"),
        ("Threat-Specific", "Beaconing analysis, crypto mining, SMB ransomware, lateral movement, exfil, scanning"),
    ]
    y = Inches(5.15)
    for cat, tables in categories:
        _add_rect(sl, Inches(0.3), y, Inches(2.2), Inches(0.28), fill_color=BLUE_MID)
        _add_textbox(sl, Inches(0.4), y + Inches(0.02), Inches(2.0), Inches(0.24),
                     cat, font_size=Pt(8.5), bold=True, color=WHITE)
        _add_textbox(sl, Inches(2.65), y + Inches(0.02), Inches(10.0), Inches(0.24),
                     tables, font_size=Pt(8.5), color=MID_GREY)
        y += Inches(0.3)

    _footer(sl)
    _speaker_note(sl,
        "The query generator reads the case IOCs and produces ready-to-paste "
        "queries for three SIEM platforms. It also detects threat patterns in "
        "the investigation report — like beaconing or lateral movement — and "
        "generates threat-specific hunt queries automatically. When a Sentinel "
        "schema registry is available, queries are scoped to confirmed tables.")


# ── Slide 8: PER Environment Hunt ────────────────────────────────────────────

def slide_per_environment(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, NAVY)
    _slide_header(sl, "Hunting in PER Environment",
                  "Queries auto-scoped to example-client (PER) Sentinel workspace — 57 tables available")

    # Workspace info box
    _add_rect(sl, Inches(0.3), Inches(1.0), Inches(5.5), Inches(1.3),
              fill_color=BLUE_MID)
    _add_textbox(sl, Inches(0.5), Inches(1.05), Inches(5.1), Inches(0.35),
                 "PER WORKSPACE", font_size=Pt(13), bold=True, color=BLUE_LIGHT)
    workspace_lines = [
        "Workspace ID:  00000000-0000-0000-0000-000000000000",
        "Tables:        57 (auto-discovered via Sentinel schema)",
        "Schema:        config/sentinel_tables.json",
        "CLI:           az monitor log-analytics query -w <ID>",
    ]
    y = Inches(1.42)
    for line in workspace_lines:
        _add_textbox(sl, Inches(0.5), y, Inches(5.1), Inches(0.25),
                     line, font_size=Pt(9.5), color=LIGHT_GREY)
        y += Inches(0.25)

    # Key hunt tables available in PER
    _add_textbox(sl, Inches(6.1), Inches(1.0), Inches(6.8), Inches(0.35),
                 "KEY HUNT TABLES IN PER", font_size=Pt(13), bold=True, color=BLUE_LIGHT)

    per_tables = [
        ("DeviceNetworkEvents",    "Network connections — IP, domain, URL hunts"),
        ("DeviceFileEvents",       "File system — SHA256/MD5 hash hunts"),
        ("DeviceProcessEvents",    "Process creation — command line + hash hunts"),
        ("DeviceImageLoadEvents",  "DLL/module loads — implant detection"),
        ("DeviceLogonEvents",      "Device logons — lateral movement"),
        ("DeviceRegistryEvents",   "Registry changes — persistence detection"),
        ("IdentityLogonEvents",    "Authentication — brute force / credential spray"),
        ("IdentityQueryEvents",    "AD queries — reconnaissance detection"),
        ("SecurityEvent",          "Windows security events (4624/4625/4648)"),
        ("EmailEvents",            "Email delivery — phishing hunts"),
        ("UrlClickEvents",         "Safe Links clicks — click-through tracking"),
        ("CommonSecurityLog",      "CEF syslog — firewall/proxy logs"),
        ("Syslog",                 "Linux/Unix syslog"),
        ("AlertEvidence",          "MDE alert evidence — SHA256 correlation"),
    ]

    col_widths = [Inches(2.8), Inches(3.9)]
    hdr_y = Inches(1.4)
    x0 = Inches(6.1)

    for j, (hdr, w) in enumerate(zip(["Table", "Hunt Capability"], col_widths)):
        _add_rect(sl, x0 + sum(col_widths[:j]), hdr_y, w, Inches(0.28),
                  fill_color=BLUE_LIGHT)
        _add_textbox(sl, x0 + sum(col_widths[:j]) + Inches(0.05), hdr_y + Inches(0.02),
                     w - Inches(0.1), Inches(0.24),
                     hdr, font_size=Pt(9), bold=True, color=NAVY)

    y = hdr_y + Inches(0.28)
    for i, (table, desc) in enumerate(per_tables):
        bg = DARK_BG if i % 2 == 0 else DARKER_BG
        cx = x0
        for val, w in zip([table, desc], col_widths):
            _add_rect(sl, cx, y, w, Inches(0.3), fill_color=bg)
            _add_textbox(sl, cx + Inches(0.05), y + Inches(0.02),
                         w - Inches(0.1), Inches(0.26),
                         val, font_size=Pt(8.5), color=LIGHT_GREY)
            cx += w
        y += Inches(0.3)

    # Live KQL example scoped to PER
    _add_textbox(sl, Inches(0.3), Inches(2.5), Inches(5.5), Inches(0.35),
                 "EXAMPLE: SCOPED HUNT COMMAND", font_size=Pt(12), bold=True, color=BLUE_LIGHT)

    cmd = (
        'python3 socai.py queries --case IV_CASE_003 \\\n'
        '    --platforms kql \\\n'
        '    --tables DeviceNetworkEvents \\\n'
        '            DeviceFileEvents \\\n'
        '            DeviceProcessEvents \\\n'
        '            DeviceImageLoadEvents \\\n'
        '            IdentityLogonEvents \\\n'
        '            SecurityEvent'
    )
    _code_block(sl, Inches(0.3), Inches(2.9), Inches(5.5), Inches(1.7), cmd, font_size=Pt(9))

    # Generated KQL output for PER
    _add_textbox(sl, Inches(0.3), Inches(4.75), Inches(5.5), Inches(0.35),
                 "GENERATED KQL — HELLOMEMORY DLL HUNT", font_size=Pt(12), bold=True, color=AMBER)

    kql = (
        '// DeviceImageLoadEvents — DLL loads [SHA256]\n'
        '// Hunt for HelloMemory.dll being loaded\n'
        'DeviceImageLoadEvents\n'
        '| where SHA256 in (\n'
        '    "160dc775c7fccd4ac45c250944f5d3cd\\\n'
        '     4f3f8414ba2587742f0fb85733fa48c7")\n'
        '| project TimeGenerated, DeviceName,\n'
        '    FileName, FolderPath, SHA256,\n'
        '    InitiatingProcessFileName\n'
        '| order by TimeGenerated asc\n'
        '\n'
        '// DeviceRegistryEvents — persistence\n'
        'DeviceRegistryEvents\n'
        '| where RegistryKey has\n'
        '    "CurrentVersion\\\\Run"\n'
        '    and ActionType == "RegistryValueSet"\n'
        '| project TimeGenerated, DeviceName,\n'
        '    RegistryKey, RegistryValueName,\n'
        '    RegistryValueData, ActionType\n'
        '| order by TimeGenerated asc'
    )
    _code_block(sl, Inches(0.3), Inches(5.15), Inches(5.5), Inches(1.8),
                kql, font_size=Pt(8), color=AMBER)

    # Callout: schema-aware scoping
    _add_rect(sl, Inches(6.1), Inches(5.75), Inches(6.8), Inches(1.0),
              fill_color=DARKER_BG, line_color=TEAL, line_width=Pt(1.5))
    _add_textbox(sl, Inches(6.3), Inches(5.8), Inches(6.4), Inches(0.35),
                 "SCHEMA-AWARE QUERY SCOPING", font_size=Pt(12), bold=True, color=TEAL)
    _add_textbox(sl, Inches(6.3), Inches(6.15), Inches(6.4), Inches(0.55),
                 "SOCAI auto-discovers Sentinel workspace schemas via "
                 "scripts/discover_sentinel_schemas.py. When --tables is provided, "
                 "queries are scoped to only confirmed tables in the target workspace — no queries "
                 "for tables that don't exist in the target environment. "
                 "Workspace catalogued: example-client.",
                 font_size=Pt(9.5), color=LIGHT_GREY)

    _footer(sl)
    _speaker_note(sl,
        "PER is the example-client Sentinel workspace with 57 tables. "
        "When you scope queries to PER's confirmed tables, SOCAI only generates "
        "KQL for tables that actually exist — no wasted time running queries "
        "against non-existent tables. The HelloMemory case (IV_CASE_003) generates "
        "targeted hunts for the implant DLL hash in DeviceImageLoadEvents, "
        "C2 connections in DeviceNetworkEvents, and persistence in "
        "DeviceRegistryEvents. These queries can be pasted directly into "
        "the Sentinel query workspace or automated via the YAML export.")


# ── Slide 9: YAML Export ────────────────────────────────────────────────────

def slide_yaml_export(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, NAVY)
    _slide_header(sl, "YAML Query Export",
                  "Machine-readable structured output for SIEM automation and CI/CD integration")

    # Left: YAML example
    yaml_example = (
        'metadata:\n'
        '  case_id: IV_CASE_003\n'
        '  generated: "2026-02-23 23:44 UTC"\n'
        '  title: "Suspicious PowerShell script"\n'
        '  severity: HIGH\n'
        '  ioc_summary: "1 ipv4, 3 domain, 5 url,\n'
        '    3 md5, 3 sha1, 5 sha256"\n'
        '  threat_patterns:\n'
        '    - beaconing\n'
        '  platforms:\n'
        '    - kql\n'
        '    - splunk\n'
        '    - logscale\n'
        '  tables: null   # or scoped list\n'
        '\n'
        'queries:\n'
        '  - platform: kql\n'
        '    category: ipv4\n'
        '    table: DeviceNetworkEvents\n'
        '    description: Network connections\n'
        '    query: |\n'
        '      let suspect_ips = dynamic(\n'
        '        ["89.42.218.223"]);\n'
        '      DeviceNetworkEvents\n'
        '      | where RemoteIP in (suspect_ips)\n'
        '          or LocalIP in (suspect_ips)\n'
        '      | project TimeGenerated,\n'
        '          DeviceName, RemoteIP,\n'
        '          RemotePort, ActionType\n'
        '      | order by TimeGenerated asc\n'
        '\n'
        '  - platform: splunk\n'
        '    category: ipv4\n'
        '    description: IPv4 Lookups\n'
        '    query: |\n'
        '      index=* (src_ip IN (\n'
        '        89.42.218.223) ...)\n'
        '\n'
        '  - platform: logscale\n'
        '    category: ipv4\n'
        '    description: "IPv4 - Network"\n'
        '    query: |\n'
        '      #event_simpleName =\n'
        '        NetworkConnectIP4 ...'
    )

    _code_block(sl, Inches(0.3), Inches(1.0), Inches(6.0), Inches(5.8),
                yaml_example, font_size=Pt(8.5))

    # Right: Benefits and integration
    _add_textbox(sl, Inches(6.6), Inches(1.0), Inches(6.4), Inches(0.35),
                 "WHY YAML?", font_size=Pt(15), bold=True, color=BLUE_LIGHT)

    benefits = [
        ("Machine-Readable",    "Parse with any language — Python, Go, PowerShell"),
        ("SIEM Automation",     "Ingest into Sentinel Analytic Rules or Splunk saved searches"),
        ("CI/CD Pipelines",     "Auto-deploy hunt queries as part of detection engineering"),
        ("Version Control",     "Track query evolution alongside case artefacts in git"),
        ("Cross-Platform",      "Single file covers KQL + Splunk + LogScale"),
    ]

    y = Inches(1.5)
    for title, desc in benefits:
        _add_rect(sl, Inches(6.6), y, Inches(6.2), Inches(0.72),
                  fill_color=DARK_BG, line_color=BLUE_MID, line_width=Pt(0.75))
        _add_textbox(sl, Inches(6.8), y + Inches(0.04), Inches(5.8), Inches(0.3),
                     title, font_size=Pt(12), bold=True, color=WHITE)
        _add_textbox(sl, Inches(6.8), y + Inches(0.35), Inches(5.8), Inches(0.3),
                     desc, font_size=Pt(10), color=LIGHT_GREY)
        y += Inches(0.82)

    # Generation command
    _add_textbox(sl, Inches(6.6), Inches(5.65), Inches(6.2), Inches(0.35),
                 "GENERATE WITH ONE COMMAND", font_size=Pt(12), bold=True, color=BLUE_LIGHT)
    cmd = (
        'python3 socai.py queries --case IV_CASE_003 \\\n'
        '    --platforms kql splunk logscale \\\n'
        '    --tables DeviceNetworkEvents \\\n'
        '            IdentityLogonEvents \\\n'
        '            SecurityEvent Syslog'
    )
    _code_block(sl, Inches(6.6), Inches(6.0), Inches(6.2), Inches(0.85),
                cmd, font_size=Pt(9))

    _footer(sl)
    _speaker_note(sl,
        "The YAML export is the bridge between investigation and response. "
        "It contains structured metadata plus every generated query with its "
        "platform, category, target table, and description. This can be "
        "consumed by SIEM automation, detection engineering pipelines, or "
        "simply version-controlled alongside the case for audit purposes.")


# ── Slide 9: Full Capabilities Matrix ───────────────────────────────────────

def slide_capabilities_matrix(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, NAVY)
    _slide_header(sl, "Full Capabilities Matrix",
                  "Everything SOCAI does — from ingestion to executive reporting")

    capabilities = [
        ("INVESTIGATION",  BLUE_MID, [
            "Multi-input: URL, domain, file (ZIP auto-extract), .eml, logs",
            "Playwright web capture with SPA handling + Cloudflare detection",
            "Recursive URL crawl (configurable depth + max URLs per level)",
            "Static file analysis: PE headers, strings, entropy, embedded IOCs",
            "Email parse: SPF/DKIM/DMARC, spoofing detect, homoglyph check",
        ]),
        ("INTELLIGENCE", TEAL, [
            "15+ enrichment providers (parallel, cached, fault-isolated)",
            "Composite verdict scoring with confidence classification",
            "Cross-case IOC index with recurring infrastructure detection",
            "Campaign clustering via Union-Find on shared IOC adjacency",
            "Sandbox query + optional live detonation (Any.Run, Joe, HA)",
        ]),
        ("DETECTION", AMBER, [
            "Hunt queries for KQL, Splunk SPL, LogScale (auto-generated)",
            "YAML structured export for SIEM automation / CI/CD",
            "Threat pattern detection: beaconing, crypto, ransomware, lateral",
            "Sentinel schema registry: auto-discovered table/field mapping",
            "Behavioural anomaly detection: temporal, travel, brute force",
        ]),
        ("REPORTING", PURPLE, [
            "Markdown investigation report with 11 configurable sections",
            "PPTX executive briefing deck (management-ready)",
            "LLM security architecture review (MITRE ATT&CK, remediation)",
            "Weekly rollup across all cases",
            "MCP server for Claude Desktop integration",
        ]),
    ]

    cw = Inches(6.25)
    ch = Inches(2.55)
    gx = Inches(0.2)
    gy = Inches(0.12)

    for i, (title, color, items) in enumerate(capabilities):
        col = i % 2
        row = i // 2
        cx = Inches(0.25) + col * (cw + gx)
        cy = Inches(1.0) + row * (ch + gy)

        _add_rect(sl, cx, cy, cw, Inches(0.4), fill_color=color)
        _add_textbox(sl, cx + Inches(0.1), cy + Inches(0.04), cw - Inches(0.2),
                     Inches(0.32), title, font_size=Pt(13), bold=True, color=WHITE)

        item_y = cy + Inches(0.42)
        for item in items:
            _add_textbox(sl, cx + Inches(0.15), item_y, cw - Inches(0.3),
                         Inches(0.35), f"  {item}",
                         font_size=Pt(9.5), color=LIGHT_GREY)
            item_y += Inches(0.35)

    # Stats bar
    stats = [
        ("Enrichment Providers", "15+"),
        ("SIEM Platforms", "3"),
        ("Query Categories", "7"),
        ("Anomaly Detectors", "6"),
        ("Artefact Types", "18"),
    ]
    sx = Inches(0.25)
    sy = Inches(6.3)
    sw = Inches(2.4)
    sh = Inches(0.55)
    for i, (label, val) in enumerate(stats):
        _add_label_box(sl, sx + i * (sw + Inches(0.15)), sy, sw, sh,
                       label, val, box_color=BLUE_MID,
                       label_size=Pt(8), val_size=Pt(16))

    _footer(sl)
    _speaker_note(sl,
        "This is the full capability set. The platform handles the entire "
        "investigation lifecycle from initial alert through enrichment, "
        "detection engineering, and executive reporting. Everything is "
        "file-system based — no database required — making it portable "
        "and auditable.")


# ── Slide 10: Closing ───────────────────────────────────────────────────────

def slide_closing(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, NAVY)
    _add_rect(sl, 0, 0, SLIDE_W, Inches(0.12), fill_color=BLUE_LIGHT)

    _add_textbox(sl, Inches(0.5), Inches(1.5), Inches(12.5), Inches(1.0),
                 "SOCAI", font_size=Pt(56), bold=True, color=WHITE)

    _add_textbox(sl, Inches(0.5), Inches(2.5), Inches(12.5), Inches(0.5),
                 "From alert to actionable intelligence — automated.",
                 font_size=Pt(22), color=BLUE_LIGHT)

    # Key takeaways
    takeaways = [
        "15-step orchestrated pipeline with parallel execution and fault isolation",
        "20+ IOCs extracted automatically from 7 artefact sources in IV_CASE_003",
        "15+ enrichment providers queried in parallel with cross-run caching",
        "Composite verdict scoring with multi-provider consensus",
        "Hunt queries generated for KQL, Splunk, and LogScale simultaneously",
        "YAML export bridges investigation and SIEM automation",
        "Zero database dependency — filesystem-only, git-friendly, fully auditable",
    ]
    y = Inches(3.3)
    for t in takeaways:
        _add_textbox(sl, Inches(0.7), y, Inches(12.0), Inches(0.4),
                     f"  {t}", font_size=Pt(13), color=LIGHT_GREY)
        y += Inches(0.45)

    _add_rect(sl, 0, Inches(7.2), SLIDE_W, Inches(0.3), fill_color=BLUE_MID)
    _add_textbox(sl, Inches(0.2), Inches(7.22), Inches(6), Inches(0.25),
                 "SOCAI  |  Automated SOC Investigation Platform",
                 font_size=Pt(9), color=LIGHT_GREY)
    _speaker_note(sl,
        "Key message: SOCAI turns hours of manual analyst work into minutes "
        "of automated investigation. The IV_CASE_003 case study showed the full "
        "pipeline from PowerShell script submission through IOC extraction, "
        "enrichment, verdict scoring, and hunt query generation — all with "
        "structured YAML output ready for SIEM automation.")


# ── Main ──────────────────────────────────────────────────────────────────────

def generate_capabilities_pptx() -> str:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_pipeline(prs)
    slide_threat_actor(prs)
    slide_ioc_collection(prs)
    slide_enrichment(prs)
    slide_verdict_scoring(prs)
    slide_hunt_queries(prs)
    slide_per_environment(prs)
    slide_yaml_export(prs)
    slide_capabilities_matrix(prs)
    slide_closing(prs)

    out_path = Path(__file__).resolve().parent.parent / "socai_capabilities_showcase.pptx"
    prs.save(str(out_path))

    print(f"[generate_capabilities_pptx] {len(prs.slides)} slides written to {out_path}")
    print(f"[generate_capabilities_pptx] Size: {out_path.stat().st_size:,} bytes")
    return str(out_path)


if __name__ == "__main__":
    generate_capabilities_pptx()
